from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide - Student Details
slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "IPL Data Analytics Project"
subtitle.text = """Student Name: [Your Full Name]
College: [Your College Name]
Course: Bachelor of Computer Applications (BCA)
Project Guide: [Guide Name]
Academic Year: 2023-2024
Submission Date: [Date]"""

# Slide 2: Introduction
slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
title = slide.shapes.title
title.text = "Introduction"
content = slide.placeholders[1].text_frame
content.text = "Project Overview:"
p = content.add_paragraph()
p.text = "- Indian Premier League (IPL) data analysis from 2008-2022"
p = content.add_paragraph()
p.text = "- Comprehensive analysis using Python, data visualization, and ML"
p = content.add_paragraph()
p.text = "- Covers 1000+ matches across 15 seasons"
p = content.add_paragraph()
p.text = "Target Audience: Cricket analysts, team managers, data enthusiasts"

# Slide 3: Problem Statement
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Problem Statement"
content = slide.placeholders[1].text_frame
content.text = "Current Challenges:"
p = content.add_paragraph()
p.text = "- Data volume & complexity (1000+ matches)"
p = content.add_paragraph()
p.text = "- Lack of centralized analytics tools"
p = content.add_paragraph()
p.text = "- Limited trend identification"
p = content.add_paragraph()
p.text = "Solution: EDA, predictive modeling, visualizations"

# Slide 4: Objective
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Objective"
content = slide.placeholders[1].text_frame
content.text = "Primary Objectives:"
p = content.add_paragraph()
p.text = "- Data loading & preprocessing"
p = content.add_paragraph()
p.text = "- Exploratory Data Analysis (EDA)"
p = content.add_paragraph()
p.text = "- Pattern discovery & statistical analysis"
p = content.add_paragraph()
p.text = "- Predictive modeling (Regression & Classification)"
p = content.add_paragraph()
p.text = "- Comprehensive visualizations"

# Slide 5: Project Scope
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Project Scope"
content = slide.placeholders[1].text_frame
content.text = "In-Scope:"
p = content.add_paragraph()
p.text = "- IPL matches 2008-2022 (15 seasons)"
p = content.add_paragraph()
p.text = "- Match-level & ball-by-ball data"
p = content.add_paragraph()
p.text = "- EDA, statistical analysis, ML modeling"
p = content.add_paragraph()
p.text = "Out-of-Scope: Real-time prediction, player-level metrics"

# Slide 6: System Requirements
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "System Requirements"
content = slide.placeholders[1].text_frame
content.text = "Hardware:"
p = content.add_paragraph()
p.text = "- Processor: Intel Core i5+"
p = content.add_paragraph()
p.text = "- RAM: 8GB+"
p = content.add_paragraph()
p.text = "Software:"
p = content.add_paragraph()
p.text = "- Python 3.7+, Jupyter Notebook"
p = content.add_paragraph()
p.text = "- Libraries: pandas, numpy, matplotlib, seaborn, scikit-learn"

# Slide 7: Overview of Python
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Overview of Python"
content = slide.placeholders[1].text_frame
content.text = "Why Python for Data Analytics:"
p = content.add_paragraph()
p.text = "- Rich ecosystem (pandas, numpy, matplotlib)"
p = content.add_paragraph()
p.text = "- Easy to learn & versatile"
p = content.add_paragraph()
p.text = "- Excellent for data manipulation & visualization"
p = content.add_paragraph()
p.text = "Key Libraries: pandas (data), numpy (numerical), matplotlib/seaborn (viz), scikit-learn (ML)"

# Slide 8: Overview of SQL
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Overview of SQL"
content = slide.placeholders[1].text_frame
content.text = "Role in Project:"
p = content.add_paragraph()
p.text = "- Conceptual understanding of database design"
p = content.add_paragraph()
p.text = "- SQL queries for data retrieval (conceptual)"
p = content.add_paragraph()
p.text = "Tables: IPL_Matches, Ball_by_Ball, Teams, Venues"

# Slide 9: Project Module
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Project Module"
content = slide.placeholders[1].text_frame
content.text = "Modules:"
p = content.add_paragraph()
p.text = "- Data Loading & Preprocessing"
p = content.add_paragraph()
p.text = "- Exploratory Data Analysis"
p = content.add_paragraph()
p.text = "- Statistical Analysis"
p = content.add_paragraph()
p.text = "- Visualization"
p = content.add_paragraph()
p.text = "- Predictive Modeling"
p = content.add_paragraph()
p.text = "- Reporting"

# Slide 10: Data Flow Diagram
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Data Flow Diagram"
content = slide.placeholders[1].text_frame
content.text = "Process Flow:"
p = content.add_paragraph()
p.text = "CSV Files → Data Loading → Preprocessing → EDA → Statistical Analysis → Modeling → Visualization → Reports"

# Slide 11: Database Design
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Database Design"
content = slide.placeholders[1].text_frame
content.text = "Tables:"
p = content.add_paragraph()
p.text = "- IPL_Matches: MatchID, Season, Teams, WinningTeam, WinMargin, etc."
p = content.add_paragraph()
p.text = "- Ball_by_Ball: BallID, MatchID, Batsman, Bowler, Runs, etc."
p = content.add_paragraph()
p.text = "Normalization: 3NF, Relationships via foreign keys"

# Slide 12: Python Code
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Python Code"
content = slide.placeholders[1].text_frame
content.text = "Key Code Snippets:"
p = content.add_paragraph()
p.text = "- Data Loading: pd.read_csv()"
p = content.add_paragraph()
p.text = "- EDA: df.describe(), df.corr()"
p = content.add_paragraph()
p.text = "- Modeling: LinearRegression, RandomForestClassifier"
p = content.add_paragraph()
p.text = "- Visualization: plt.plot(), sns.heatmap()"

# Slide 13: Outputs
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Outputs"
content = slide.placeholders[1].text_frame
content.text = "Results:"
p = content.add_paragraph()
p.text = "- Dataset: 1000 matches, clean data"
p = content.add_paragraph()
p.text = "- Models: R²=0.76 (Regression), Accuracy=82% (Classification)"
p = content.add_paragraph()
p.text = "- Visualizations: 20+ charts (histograms, heatmaps, etc.)"

# Slide 14: Data Analysis (Graphs)
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Data Analysis (Graphs)"
content = slide.placeholders[1].text_frame
content.text = "Key Visualizations:"
p = content.add_paragraph()
p.text = "- Winning Margin Distribution (Histogram)"
p = content.add_paragraph()
p.text = "- Team-wise Win Distribution (Bar Chart)"
p = content.add_paragraph()
p.text = "- Correlation Heatmap"
p = content.add_paragraph()
p.text = "- Actual vs Predicted Scatter Plot"
p = content.add_paragraph()
p.text = "Insights: Increasing margins over seasons, home advantage"

# Slide 15: Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Conclusion"
content = slide.placeholders[1].text_frame
content.text = "Key Findings:"
p = content.add_paragraph()
p.text = "- Winning margins increasing over seasons"
p = content.add_paragraph()
p.text = "- Home advantage effect (10%)"
p = content.add_paragraph()
p.text = "- Models predict outcomes with 82%+ accuracy"
p = content.add_paragraph()
p.text = "Future: Real-time predictions, player analytics"

# Slide 16: Bibliography
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Bibliography"
content = slide.placeholders[1].text_frame
content.text = "References:"
p = content.add_paragraph()
p.text = "- Pandas Documentation"
p = content.add_paragraph()
p.text = "- Scikit-learn Documentation"
p = content.add_paragraph()
p.text = "- IPL Official Website"
p = content.add_paragraph()
p.text = "- Academic papers on sports analytics"

# Save the presentation
prs.save('IPL_Data_Analytics_Presentation.pptx')
print("PPT created successfully!")